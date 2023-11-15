import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from sklearn.feature_extraction.text import CountVectorizer, TfidfVectorizer
from nltk.sentiment.vader import SentimentIntensityAnalyzer
from io import BytesIO
import openpyxl
from openpyxl.utils.dataframe import dataframe_to_rows
from openpyxl.drawing.image import Image
from pptx import Presentation
from pptx.util import Inches
import datetime
import tempfile
import os
import shutil

import constants as c

# # # # #
# temp vars

df = pd.read_excel(c.data_file)

master_pptx = c.master_pptx 
output_pptx = c.output_pptx

# # # # #

def textAnalysis(df=df):

    postings = df

    cvec = CountVectorizer(stop_words='english')# fit cvec to text column
    cvec.fit(postings['User Message'])# transform text column
    text_cvec = cvec.transform(postings['User Message'])# create dataframe with the transformed vectors
    df = pd.DataFrame(text_cvec.todense(),
                    columns = cvec.get_feature_names_out())# view results
    #print(df)

    #visualizing
    # set image size
    chart1 = plt.figure(figsize=(15, 7))
    data = df.sum().sort_values(ascending=False).head(30)
    ax = sns.barplot(x=data.values, y=data.index, hue=data.index, palette='coolwarm', legend=False)
    ax.set_xlabel('Häufigkeit')
    ax.set_ylabel('Wörter')
    ax.set_title('Die häufigsten Worte (Ohne Stop-Worte)')
    sns.despine(left=True, bottom=True)  # Entfernt den Rahmen
    #plt.show()

    # instantiate sentiment analysis
    sentiment = SentimentIntensityAnalyzer()# create an empty list for sentiment scores
    sentiment_list = []# loop through each tweet and calculate overall sentiment score, then append to the list above
    for i in postings['User Message']:
        score = sentiment.polarity_scores(i)['compound']
        sentiment_list.append(score)# create a new column in postings dataframe that has the sentiment scores
    postings['sentiment'] = sentiment_list# view results
    #print(postings.head())

    # create new dataframe with positive sentiment postings
    positives = postings[postings['sentiment'] > 0]# instantiate CountVectorizer with stop_words parameter
    cvec_pos = CountVectorizer(stop_words='english')# fit cvec to text column
    cvec_pos.fit(positives['User Message'])# transform text column
    text_cvec_pos = cvec_pos.transform(positives['User Message'])# create dataframe with the transformed vectors
    df_pos = pd.DataFrame(text_cvec_pos.todense(),
                    columns = cvec_pos.get_feature_names_out())# set image size

    chart2 = plt.figure(figsize=(15, 7))
    data = df_pos.sum().sort_values(ascending=False).head(30)
    sns.barplot(x=data.values, y=data.index, hue=data.index, palette='coolwarm', legend=False)
    sns.despine(left=True, bottom=True)  # Entfernt den Rahmen
    plt.xlabel('Häufigkeit')
    plt.ylabel('Wörter')
    plt.title('Die häufigsten Worte in positiven Postings')

    # create new dataframe with negative sentiment postings
    negatives = postings[postings['sentiment'] < 0]# instantiate CountVectorizer with stop_words parameter
    cvec_neg = CountVectorizer(stop_words='english')# fit cvec to text column
    cvec_neg.fit(negatives['User Message'])# transform text column
    text_cvec_neg = cvec_neg.transform(negatives['User Message'])# create dataframe with the transformed vectors
    df_neg = pd.DataFrame(text_cvec_neg.todense(),
                    columns = cvec_neg.get_feature_names_out())# set image size
    chart3 = plt.figure(figsize=(15, 7))
    data = df_neg.sum().sort_values(ascending=False).head(30)
    ax = sns.barplot(x=data.values, y=data.index, hue=data.index, palette='coolwarm', legend=False)
    ax.set_xlabel('Häufigkeit')
    ax.set_ylabel('Wörter')
    ax.set_title('Die häufigsten Worte in negativen Postings')
    sns.despine(left=True, bottom=True)  # Entfernt den Rahmen
    #plt.show()

    return chart1, chart2, chart3


def createExport(passToPresentation):
    shutil.copy(master_pptx, output_pptx)

    prs = Presentation(output_pptx)

    buffer_list = [BytesIO() for _ in range(3)]
    images = [passToPresentation[i] for i in range(3)]

    for i, image in enumerate(images):
        image.savefig(buffer_list[i], format="png")
        buffer_list[i].seek(0)

        # Füge eine neue Folie hinzu
        slide_layout = prs.slide_layouts[0]  # Hier wird der Layout-Index für Folie 2 verwendet
        slide = prs.slides.add_slide(slide_layout)
        
        # Fügen Sie das Bild auf der Folie ein
        left = Inches(1)  # Anpassen der Position nach Bedarf
        top = Inches(1)   # Anpassen der Position nach Bedarf
        pic = slide.shapes.add_picture(buffer_list[i], left, top, height=Inches(4), width=Inches(8))

    # Speichere die aktualisierte Präsentation
    prs.save(output_pptx)

    print("PowerPoint presentation created and saved successfully.")

